"""
Track D1 — export breadth: the same canonical ReportObject rendered as a
branded PPTX deck and a responsive email-ready HTML page, behind one shared
interface so Google Slides (needs an OAuth connector not yet wired up) can
slot in later without callers changing shape.

Same rule as every other renderer in this app: each export function's only
input is the object. No re-querying, no recompute -- a number in a slide or
an email is the same number report.html/the dashboard already showed,
because it's the same obj.metrics/obj.narrative/obj.qa.
"""
from __future__ import annotations

import base64
import io
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

from . import theme
from .html_dashboard import _kpi_cards
from .report_object import ReportObject

PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
EMAIL_HTML_CONTENT_TYPE = "text/html"


@dataclass
class ExportResult:
    format: str  # "pptx" | "email_html" | "google_slides"
    status: str  # "ok" | "unavailable"
    content: bytes | str | None = None
    content_type: str | None = None
    reason: str | None = None


def _rgb(hex_color: str | None, fallback: str) -> RGBColor:
    value = (hex_color or fallback).lstrip("#")
    return RGBColor.from_string(value.upper())


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _add_title_slide(prs: Presentation, obj: ReportObject) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = _rgb(obj.branding.get("primary_color"), theme.CATEGORICAL[0])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Emu(int(prs.slide_height * 0.35)))
    bg.fill.solid()
    bg.fill.fore_color.rgb = primary
    bg.line.fill.background()
    bg.shadow.inherit = False

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    title = obj.narrative.get("report_title") or f"{obj.branding.get('client_name', 'Client')} — Performance Report"
    tf.text = title
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor.from_string("FFFFFF")

    meta_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(9), Inches(1))
    meta = meta_box.text_frame
    meta.text = f"Prepared for {obj.branding.get('client_name', 'Client')} by {obj.branding.get('agency_name', 'Your Agency')}"
    meta.paragraphs[0].font.size = Pt(14)
    p2 = meta.add_paragraph()
    p2.text = f"Reporting period: {obj.period.label}"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor.from_string("666666")


def _add_kpi_slide(prs: Presentation, obj: ReportObject) -> None:
    cards = _kpi_cards(obj.metrics)
    if not cards:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    header.text_frame.text = "Key Metrics"
    header.text_frame.paragraphs[0].font.size = Pt(24)
    header.text_frame.paragraphs[0].font.bold = True

    accent = _rgb(obj.branding.get("accent_color"), theme.CATEGORICAL[3])
    cols = min(len(cards), 4)
    card_width = Inches(9 / cols)
    for i, card in enumerate(cards[:8]):
        col, row = i % cols, i // cols
        left = Inches(0.5) + card_width * col
        top = Inches(1.4) + Inches(1.6) * row
        box = slide.shapes.add_textbox(left, top, card_width - Inches(0.15), Inches(1.4))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = str(card["formatted"])
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = accent
        label_p = tf.add_paragraph()
        label_p.text = str(card["label"])
        label_p.font.size = Pt(11)
        label_p.font.color.rgb = RGBColor.from_string("666666")


def _add_section_slides(prs: Presentation, obj: ReportObject) -> None:
    charts_by_section: dict[str, list] = {}
    for chart in obj.charts:
        charts_by_section.setdefault(chart.section, []).append(chart)

    for i, section in enumerate(obj.narrative.get("sections", [])):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        header = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
        header.text_frame.text = section.get("heading", "")
        header.text_frame.paragraphs[0].font.size = Pt(22)
        header.text_frame.paragraphs[0].font.bold = True

        narrative_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.0))
        narrative_box.text_frame.word_wrap = True
        narrative_box.text_frame.text = section.get("narrative", "")
        narrative_box.text_frame.paragraphs[0].font.size = Pt(13)

        section_key = obj.section_order[i] if i < len(obj.section_order) else None
        section_charts = charts_by_section.get(section_key, [])
        if section_charts:
            first = section_charts[0]
            try:
                img_bytes = base64.b64decode(first.img)
                slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(1.5), Inches(2.2), width=Inches(6))
            except Exception:
                pass  # a malformed image never blocks the rest of the deck


def export_pptx(obj: ReportObject) -> ExportResult:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)

    _add_title_slide(prs, obj)
    _add_kpi_slide(prs, obj)
    _add_section_slides(prs, obj)

    buf = io.BytesIO()
    prs.save(buf)
    return ExportResult(format="pptx", status="ok", content=buf.getvalue(), content_type=PPTX_CONTENT_TYPE)


# ---------------------------------------------------------------------------
# Email-ready responsive HTML
# ---------------------------------------------------------------------------

_EMAIL_TEMPLATE = """<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{title}</title>
<style>
  body {{ margin:0; padding:0; background:#f2f2f0; font-family: Helvetica, Arial, sans-serif; }}
  .container {{ max-width: 600px; margin: 0 auto; background:#ffffff; }}
  .header {{ background:{primary_color}; color:#ffffff; padding:24px 28px; }}
  .header h1 {{ margin:0; font-size:20px; }}
  .header p {{ margin:6px 0 0; font-size:13px; opacity:0.9; }}
  .qa-badge {{ display:inline-block; margin-top:10px; padding:4px 10px; border-radius:999px;
               font-size:11px; font-weight:bold; background:rgba(255,255,255,0.2); }}
  .kpi-grid {{ padding: 20px 28px 4px; }}
  .kpi {{ display:inline-block; width:45%; margin:0 2% 16px; vertical-align:top; }}
  .kpi .value {{ font-size:20px; font-weight:bold; color:{accent_color}; }}
  .kpi .label {{ font-size:11px; color:#666; text-transform:uppercase; }}
  .summary {{ padding: 4px 28px 20px; font-size:14px; line-height:1.5; color:#222; }}
  .cta {{ display:block; margin: 0 28px 28px; padding:12px 20px; background:{primary_color};
          color:#ffffff; text-align:center; text-decoration:none; border-radius:6px; font-size:14px; }}
  @media only screen and (max-width: 480px) {{
    .kpi {{ width: 100%; margin: 0 0 16px; display:block; }}
  }}
</style>
</head>
<body>
<div class="container">
  <div class="header">
    <h1>{title}</h1>
    <p>Prepared for {client_name} &middot; {period_label}</p>
    {badge_html}
  </div>
  <div class="kpi-grid">
    {kpi_html}
  </div>
  <div class="summary">{summary}</div>
  <a class="cta" href="#">View the full report</a>
</div>
</body>
</html>"""


def export_email_html(obj: ReportObject) -> ExportResult:
    cards = _kpi_cards(obj.metrics)
    kpi_html = "".join(
        f'<div class="kpi"><div class="value">{c["formatted"]}</div><div class="label">{c["label"]}</div></div>'
        for c in cards[:6]
    )
    badge = (obj.qa or {}).get("badge")
    badge_html = f'<span class="qa-badge">QA: {badge}</span>' if badge else ""

    html = _EMAIL_TEMPLATE.format(
        title=obj.narrative.get("report_title") or "Performance Report",
        client_name=obj.branding.get("client_name", "Client"),
        period_label=obj.period.label,
        primary_color=obj.branding.get("primary_color") or theme.CATEGORICAL[0],
        accent_color=obj.branding.get("accent_color") or theme.CATEGORICAL[3],
        badge_html=badge_html,
        kpi_html=kpi_html,
        summary=obj.narrative.get("executive_summary", ""),
    )
    return ExportResult(format="email_html", status="ok", content=html, content_type=EMAIL_HTML_CONTENT_TYPE)


# ---------------------------------------------------------------------------
# Google Slides — same interface, red until a connector exists
# ---------------------------------------------------------------------------

def export_google_slides(obj: ReportObject) -> ExportResult:
    """Same (obj) -> ExportResult interface as the other two exports, so a
    caller can treat all three uniformly -- but genuinely unavailable until
    a Google OAuth connector is configured (see mission's EXTERNAL
    DEPENDENCIES note). Never silently returns a fake success."""
    return ExportResult(
        format="google_slides", status="unavailable",
        reason="Google Slides export requires a connected Google account (OAuth) — not yet configured.",
    )


def export_pbip(obj: ReportObject) -> ExportResult:
    """Track D2.1 — Power BI, zipped. build_pbip() writes a directory tree
    (TMDL tables, PBIR pages/visuals -- a real Power BI Project is dozens of
    small files, not one blob), so unlike pptx/email_html this export zips
    that tree into a single downloadable archive -- the natural single-file
    shape for "a folder" over a web download. Same never-a-fake-success rule
    as google_slides below: if build_pbip() raises, this propagates it
    rather than returning a status="ok" empty file."""
    import tempfile
    import zipfile

    with tempfile.TemporaryDirectory() as tmp:
        from . import pbip_export
        tmp_path = Path(tmp)
        pbip_export.build_pbip(obj, tmp_path)

        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for path in sorted(tmp_path.rglob("*")):
                if path.is_file():
                    zf.write(path, path.relative_to(tmp_path))
        return ExportResult(format="pbip", status="ok", content=buf.getvalue(), content_type="application/zip")


def export_report(obj: ReportObject, formats: list[str] | None = None) -> dict[str, ExportResult]:
    """formats: subset of ("pptx", "email_html", "google_slides", "pbip"); defaults to all."""
    exporters = {
        "pptx": export_pptx, "email_html": export_email_html,
        "google_slides": export_google_slides, "pbip": export_pbip,
    }
    formats = formats or list(exporters)
    return {fmt: exporters[fmt](obj) for fmt in formats if fmt in exporters}
