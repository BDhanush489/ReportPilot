"""Tests for app/exports.py (Track D1 — export breadth: PPTX, email-HTML,
Google Slides behind one interface)."""
import io

from pptx import Presentation

from app.exports import export_email_html, export_google_slides, export_pbip, export_pptx, export_report
from app.report_object import ChartRef, Period, ReportObject, SourceInfo

METRICS = {
    "analytics": {
        "totals": {"sessions": 1000, "revenue_usd": 5000.0, "conversion_rate": 3.0},
        # revenue_usd sums to totals.revenue_usd exactly -- D2.2's Power BI
        # measure reconciliation (SUM(AnalyticsByChannel[revenue_usd])) checks
        # this same real ReportObject via export_pbip, so the fixture must be
        # internally consistent, not just illustrative.
        "by_channel": [{"channel": "Organic Search", "revenue_usd": 5000.0, "conversion_rate": 3.0}],
    },
    "period_label": "2026-01-01 to 2026-06-30",
}

NARRATIVE = {
    "report_title": "Aurora Home Goods — Performance Report",
    "period_label": "2026-01-01 to 2026-06-30",
    "executive_summary": "Sessions grew to 1,000, driving $5,000 in revenue.",
    "highlights": ["Organic Search led revenue."],
    "watchouts": [],
    "sections": [
        {"heading": "Web Analytics", "narrative": "Revenue reached $5,000.", "recommendations": []},
    ],
    "next_steps": [],
}

BRANDING = {"agency_name": "Northlight", "client_name": "Aurora", "primary_color": "#2a78d6", "accent_color": "#eda100"}

QA_PASS = {
    "badge": "PASS", "failing_checks": [],
    "traceability": {"ok": True, "numbers_checked": 1, "fail": [], "warnings": []},
    "aggregation_sanity": {"ok": True, "mismatches": [], "inconclusive_sources": []},
    "unsupported_claims": {"ok": True, "claims_checked": 1, "unlinked": []},
}

# A minimal 1x1 PNG, base64-encoded -- enough for add_picture() to accept it.
_TINY_PNG_B64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
)


def _build_object(with_chart: bool = True) -> ReportObject:
    charts = []
    if with_chart:
        charts.append(ChartRef(
            id="analytics-0-revenue-by-channel", section="analytics", caption="Revenue by channel",
            img=_TINY_PNG_B64, chart_type="bar", metric_paths=["metrics.analytics.by_channel"],
            suitability_verdict="good", suitability_reason="fine",
        ))
    return ReportObject(
        report_id="test-export-1",
        period=Period(label=METRICS["period_label"]),
        sources={"analytics": SourceInfo(row_count=10, sha256="a" * 64)},
        metrics=METRICS,
        series={},
        charts=charts,
        narrative=NARRATIVE,
        qa=QA_PASS,
        branding=BRANDING,
        section_order=["analytics"],
    )


# ---------------------------------------------------------------------------
# PPTX: branded slides (title, KPI, per-section chart+narrative)
# ---------------------------------------------------------------------------

def test_pptx_export_succeeds_and_is_a_real_presentation():
    result = export_pptx(_build_object())
    assert result.status == "ok"
    assert result.content_type.endswith("presentationml.presentation")
    prs = Presentation(io.BytesIO(result.content))
    assert len(prs.slides) >= 3  # title, KPI, >=1 section


def test_pptx_title_slide_carries_branding_and_period():
    result = export_pptx(_build_object())
    prs = Presentation(io.BytesIO(result.content))
    slide0_text = " ".join(
        shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame
    )
    assert "Aurora Home Goods" in slide0_text
    assert "Northlight" in slide0_text
    assert METRICS["period_label"] in slide0_text


def test_pptx_kpi_slide_shows_the_same_numbers_as_metrics():
    result = export_pptx(_build_object())
    prs = Presentation(io.BytesIO(result.content))
    all_text = " ".join(
        shape.text_frame.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
    )
    assert "5,000" in all_text  # revenue_usd formatted via theme.format_currency
    assert "1,000" in all_text  # sessions formatted via theme.format_count


def test_pptx_section_slide_embeds_the_chart_picture():
    result = export_pptx(_build_object(with_chart=True))
    prs = Presentation(io.BytesIO(result.content))
    section_slide = prs.slides[-1]
    picture_count = sum(1 for shape in section_slide.shapes if shape.shape_type == 13)
    assert picture_count == 1


def test_pptx_section_slide_survives_missing_chart_without_crashing():
    result = export_pptx(_build_object(with_chart=False))
    assert result.status == "ok"
    prs = Presentation(io.BytesIO(result.content))
    assert len(prs.slides) >= 3


def test_pptx_survives_a_malformed_chart_image_without_crashing():
    obj = _build_object(with_chart=True)
    obj.charts[0].img = "not-valid-base64!!"
    result = export_pptx(obj)
    assert result.status == "ok"  # the rest of the deck still builds


# ---------------------------------------------------------------------------
# Email-ready responsive HTML
# ---------------------------------------------------------------------------

def test_email_html_export_succeeds():
    result = export_email_html(_build_object())
    assert result.status == "ok"
    assert result.content_type == "text/html"
    assert "<html" in result.content.lower()


def test_email_html_has_a_mobile_media_query_for_responsiveness():
    result = export_email_html(_build_object())
    assert "@media" in result.content
    assert "max-width" in result.content


def test_email_html_shows_the_same_numbers_as_metrics():
    result = export_email_html(_build_object())
    assert "5,000" in result.content
    assert "1,000" in result.content


def test_email_html_surfaces_the_qa_badge():
    result = export_email_html(_build_object())
    assert "PASS" in result.content


def test_email_html_omits_badge_markup_when_object_carries_no_qa():
    obj = _build_object()
    obj.qa = {}
    result = export_email_html(obj)
    # The .qa-badge CSS rule is always present in <style>; what must be
    # absent is the actual <span> element the badge text would render into.
    assert '<span class="qa-badge">' not in result.content


# ---------------------------------------------------------------------------
# Google Slides: same interface, genuinely red without a connector
# ---------------------------------------------------------------------------

def test_google_slides_export_is_explicitly_unavailable_not_a_fake_success():
    result = export_google_slides(_build_object())
    assert result.status == "unavailable"
    assert result.content is None
    assert "connector" in result.reason.lower() or "oauth" in result.reason.lower()


# ---------------------------------------------------------------------------
# Power BI (D2.1): the same (obj) -> ExportResult interface, but the
# content is a zip of build_pbip()'s whole directory tree, not one blob.
# ---------------------------------------------------------------------------

def test_pbip_export_succeeds_and_is_a_real_zip_of_a_real_project(tmp_path):
    import zipfile
    from io import BytesIO

    result = export_pbip(_build_object())
    assert result.status == "ok"
    assert result.content_type == "application/zip"

    zf = zipfile.ZipFile(BytesIO(result.content))
    assert zf.testzip() is None  # no corrupt member
    names = zf.namelist()
    assert any(n.endswith(".SemanticModel/definition.pbism") for n in names)
    assert any(n.endswith(".SemanticModel/definition/tables/AnalyticsByChannel.tmdl") for n in names)
    assert any(n.endswith(".Report/definition/report.json") for n in names)  # the fixture's one chart gets a page
    assert any(n.endswith(".pbip") for n in names)


def test_pbip_export_still_succeeds_with_no_charts_but_kpi_cards_remain():
    """No ChartRefs doesn't mean no Report/.pbip: the fixture still carries
    real metrics.analytics.totals, so a KPI card (Web Sessions/Web Revenue/
    Conversion Rate) still gets a page even with zero charts -- see
    pbip_export.py's ensure_page() being called from the KPI-card loop
    independently of the chart loop. Only a report with truly nothing to
    bind anywhere (no charts, no totals, no series) skips Report/.pbip
    entirely -- covered directly in test_pbip_export.py, not duplicated here."""
    import zipfile
    from io import BytesIO

    result = export_pbip(_build_object(with_chart=False))
    assert result.status == "ok"
    names = zipfile.ZipFile(BytesIO(result.content)).namelist()
    assert any(".SemanticModel" in n for n in names)
    assert any(".Report" in n for n in names)
    assert any(n.endswith(".pbip") for n in names)


# ---------------------------------------------------------------------------
# Shared interface + shared numbers across formats
# ---------------------------------------------------------------------------

def test_export_report_runs_all_four_formats_behind_one_call():
    results = export_report(_build_object())
    assert set(results.keys()) == {"pptx", "email_html", "google_slides", "pbip"}
    assert results["pptx"].status == "ok"
    assert results["email_html"].status == "ok"
    assert results["google_slides"].status == "unavailable"
    assert results["pbip"].status == "ok"


def test_export_report_can_be_scoped_to_a_subset_of_formats():
    results = export_report(_build_object(), formats=["email_html"])
    assert set(results.keys()) == {"email_html"}


def test_pptx_and_email_html_show_identical_revenue_figure():
    """All exports share one branding/token source; numbers identical
    across formats -- both are formatted via the same theme.format_currency
    call inside the same _kpi_cards(), not two independent f-strings that
    could drift."""
    obj = _build_object()
    pptx_result = export_pptx(obj)
    email_result = export_email_html(obj)

    prs = Presentation(io.BytesIO(pptx_result.content))
    pptx_text = " ".join(
        shape.text_frame.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
    )
    assert "5,000" in pptx_text
    assert "5,000" in email_result.content
